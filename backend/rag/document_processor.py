import os
from pathlib import Path
from typing import Optional

from pypdf import PdfReader
from pypdf.errors import FileNotDecryptedError, WrongPasswordError
from docx import Document
from pptx import Presentation
from PIL import Image
import pytesseract
pytesseract.pytesseract.tesseract_cmd = r'C:/Program Files/Tesseract-OCR/tesseract.exe'
from pdf2image import convert_from_path


SUPPORTED_EXTENSIONS = [
    "pdf",
    "docx",
    "txt",
    "md",
    "pptx",
    "jpg",
    "jpeg",
    "png",
]

IMAGE_EXTENSIONS = ["jpg", "jpeg", "png"]


def get_file_extension(file_path: str) -> str:
    return Path(file_path).suffix.lower().replace(".", "")


_OCR_DEFAULTS = {"lang": "eng", "auto_ocr": True}


def _resolve_ocr_settings(firm) -> dict:
    """
    Per-firm OCR settings (Settings > AI Configuration > OCR Configuration):
    the Tesseract language code to pass to image_to_string, and whether OCR
    is attempted at all as a fallback for scanned PDFs/images. firm=None
    (or no settings ever saved) means the defaults this pipeline has always
    used - lang="eng", auto-OCR on - so every existing caller that doesn't
    pass firm keeps behaving exactly as before.
    """
    if firm is None:
        return dict(_OCR_DEFAULTS)

    from accounts.models import FirmSettings

    try:
        settings_obj = FirmSettings.objects.get(firm=firm)
    except FirmSettings.DoesNotExist:
        return dict(_OCR_DEFAULTS)

    stored = settings_obj.data.get("ocr") or {}
    return {
        "lang": stored.get("lang") or _OCR_DEFAULTS["lang"],
        "auto_ocr": stored.get("auto_ocr", _OCR_DEFAULTS["auto_ocr"]),
    }


def extract_text_from_pdf(file_path: str, max_chars: Optional[int] = None) -> str:
    """
    Extract text from normal text-based PDF.

    max_chars stops extraction early once enough text has been collected -
    reproduced live: a 33MB/thousands-of-pages PDF took 171s and produced
    23.6M characters of text for a caller (document summarize/risk/
    compliance/entities/compare) that only ever uses the first
    MAX_DOCUMENT_CHARS (12000) of it before truncating for the LLM call -
    the other 99.9% of that 171s was pure waste, and synchronous within
    the HTTP request, which is exactly what caused the dev proxy to give
    up and reset the connection ("socket hang up") long before Django
    could respond. None (the default) means no cap - used by the RAG
    upload/chunking pipeline, which genuinely needs the whole document.
    """

    # An encrypted PDF doesn't raise when PdfReader(file_path) is called, or
    # even when .pages is assigned - reproduced live: the actual raise only
    # happens once the page LIST is iterated (pypdf resolves /Pages lazily,
    # inside enumerate()'s __iter__/__len__ call), so the whole read - not
    # just reader construction or a single page's extract_text() - has to
    # be inside this try block to actually catch it.
    try:
        reader = PdfReader(file_path)
        text_parts = []
        total_chars = 0

        for page_number, page in enumerate(reader.pages, start=1):
            page_text = page.extract_text()

            if page_text:
                text_parts.append(
                    f"\n\n--- Page {page_number} ---\n{page_text}"
                )
                total_chars += len(page_text)

            if max_chars is not None and total_chars >= max_chars:
                break
    except (FileNotDecryptedError, WrongPasswordError):
        raise ValueError(
            "This PDF is password-protected. Please remove the password and upload it again."
        )

    return "\n".join(text_parts).strip()


def extract_text_from_scanned_pdf(file_path: str, max_chars: Optional[int] = None, firm=None) -> str:
    ocr_settings = _resolve_ocr_settings(firm)
    if not ocr_settings["auto_ocr"]:
        return ""

    pages = convert_from_path(file_path)
    text_parts = []
    total_chars = 0

    for page_number, image in enumerate(pages, start=1):
        page_text = pytesseract.image_to_string(image, lang=ocr_settings["lang"])

        if page_text.strip():
            text_parts.append(
                f"\n\n--- OCR Page {page_number} ---\n{page_text.strip()}"
            )
            total_chars += len(page_text)

        if max_chars is not None and total_chars >= max_chars:
            break

    return "\n".join(text_parts).strip()


def extract_text_from_image(file_path: str, firm=None) -> str:
    """
    OCR a photo/scan of a document (JPG/PNG) using the same Tesseract
    pipeline used for scanned PDFs.
    """
    ocr_settings = _resolve_ocr_settings(firm)
    if not ocr_settings["auto_ocr"]:
        return ""

    image = Image.open(file_path)
    text = pytesseract.image_to_string(image, lang=ocr_settings["lang"])

    return text.strip()


def extract_text_from_docx(file_path: str, max_chars: Optional[int] = None) -> str:

    document = Document(file_path)
    text_parts = []
    total_chars = 0

    for paragraph in document.paragraphs:
        text = paragraph.text.strip()

        if text:
            text_parts.append(text)
            total_chars += len(text)

        if max_chars is not None and total_chars >= max_chars:
            break

    return "\n".join(text_parts).strip()


def extract_text_from_txt_or_md(file_path: str, max_chars: Optional[int] = None) -> str:
    """
    Extract text from TXT or Markdown file.
    """

    encodings = ["utf-8", "utf-16", "latin-1"]

    for encoding in encodings:
        try:
            with open(file_path, "r", encoding=encoding) as file:
                text = file.read(max_chars) if max_chars is not None else file.read()
                return text.strip()
        except UnicodeDecodeError:
            continue

    raise ValueError("Unable to read text file encoding.")


def extract_text_from_pptx(file_path: str, max_chars: Optional[int] = None) -> str:
    """
    Extract text from PPTX slides.
    """

    presentation = Presentation(file_path)
    text_parts = []
    total_chars = 0

    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_text_parts = []

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()

                if text:
                    slide_text_parts.append(text)

        if slide_text_parts:
            slide_text = "\n".join(slide_text_parts)
            text_parts.append(
                f"\n\n--- Slide {slide_number} ---\n{slide_text}"
            )
            total_chars += len(slide_text)

        if max_chars is not None and total_chars >= max_chars:
            break

    return "\n".join(text_parts).strip()


def extract_text_from_document(
    file_path: str,
    document_type: Optional[str] = None,
    max_chars: Optional[int] = None,
    firm=None,
) -> str:
    """
    Universal document extractor.

    Supports:
    - PDF
    - scanned PDF
    - DOCX
    - TXT
    - MD
    - PPTX

    max_chars: stop extraction once this many characters have been
    collected, instead of reading the whole document - pass this from a
    caller that will truncate/cap the text anyway (document
    summarize/risk/compliance/entities/compare), never from the RAG
    upload/chunking pipeline, which needs the complete document text to
    be fully searchable.

    firm: whose OCR settings (Settings > AI Configuration > OCR
    Configuration - language, auto-OCR on/off) apply when a scanned PDF or
    image needs OCR. None uses this pipeline's original defaults.
    """

    if not os.path.exists(file_path):
        raise FileNotFoundError("Document file not found.")

    extension = document_type or get_file_extension(file_path)
    extension = extension.lower()

    if extension == "pdf":
        text = extract_text_from_pdf(file_path, max_chars=max_chars)

        if not text:
            text = extract_text_from_scanned_pdf(file_path, max_chars=max_chars, firm=firm)

    elif extension == "scanned_pdf":
        text = extract_text_from_scanned_pdf(file_path, max_chars=max_chars, firm=firm)

    elif extension in IMAGE_EXTENSIONS:
        text = extract_text_from_image(file_path, firm=firm)

    elif extension == "docx":
        text = extract_text_from_docx(file_path, max_chars=max_chars)

    elif extension in ["txt", "md"]:
        text = extract_text_from_txt_or_md(file_path, max_chars=max_chars)

    elif extension == "pptx":
        text = extract_text_from_pptx(file_path, max_chars=max_chars)

    else:
        raise ValueError(
            f"Unsupported document type: {extension}. "
            "Supported types are: pdf, docx, txt, md, pptx"
        )

    if not text:
        raise ValueError(
            "No text could be extracted from this document."
        )

    return text
